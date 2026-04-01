import streamlit as st
import chromadb
from sentence_transformers import SentenceTransformer
from langchain_groq import ChatGroq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation
from langgraph.graph import StateGraph, START, END
from typing import TypedDict
import re

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle, Image
from reportlab.lib import colors
from reportlab.lib.units import cm
import urllib.request
import tempfile
import os

# ─────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────

st.set_page_config(
    page_title="Vignan's Question Paper Generator",
    page_icon="🎓",
    layout="wide"
)

# ─────────────────────────────────────────
# GLOBAL CSS
# ─────────────────────────────────────────

st.markdown("""
<style>
/* ── Google Font ── */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700;800&display=swap');

html, body, [class*="css"] {
    font-family: 'Inter', sans-serif;
}

/* ── Login page ── */
.login-wrapper {
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    min-height: 80vh;
}
.login-card {
    background: #ffffff;
    border-radius: 20px;
    padding: 44px 48px 36px 48px;
    box-shadow: 0 8px 40px rgba(26,60,110,0.13);
    width: 100%;
    max-width: 420px;
    border-top: 5px solid #1a3c6e;
}
.login-logo {
    text-align: center;
    margin-bottom: 6px;
}
.login-title {
    font-size: 22px;
    font-weight: 800;
    color: #1a3c6e;
    text-align: center;
    margin-bottom: 2px;
}
.login-sub {
    font-size: 13px;
    color: #777;
    text-align: center;
    margin-bottom: 26px;
}

/* ── Main app title bar ── */
.title-bar {
    display: flex;
    align-items: center;
    justify-content: space-between;
    background: #ffffff;
    border-radius: 14px;
    padding: 18px 28px;
    margin-bottom: 24px;
    border: 1px solid #d0dce8;
    box-shadow: 0 2px 12px rgba(26,60,110,0.07);
}
.title-left {
    display: flex;
    align-items: center;
    gap: 14px;
}
.title-icon {
    width: 46px; height: 46px;
    background: #1a3c6e;
    border-radius: 12px;
    display: flex; align-items: center; justify-content: center;
    font-size: 22px;
}
.title-text h1 {
    font-size: 22px;
    font-weight: 800;
    color: #1a3c6e;
    margin: 0; padding: 0;
    line-height: 1.2;
}
.title-text p {
    font-size: 13px;
    color: #6b7280;
    margin: 2px 0 0 0;
}
.title-right {
    display: flex;
    align-items: center;
    gap: 10px;
}
.faculty-chip {
    background: #eef2fb;
    color: #1a3c6e;
    font-size: 13px;
    font-weight: 600;
    padding: 6px 16px;
    border-radius: 20px;
    border: 1px solid #c7d5ee;
}

/* ── Cards ── */
.card {
    background: #f5f8fc;
    border-radius: 14px;
    padding: 22px 26px;
    border: 1px solid #d0dce8;
    margin-bottom: 22px;
}
.step-header {
    font-size: 17px;
    font-weight: 700;
    color: #1a3c6e;
    border-left: 4px solid #1a3c6e;
    padding-left: 10px;
    margin-bottom: 14px;
}

/* ── Sidebar ── */
.sidebar-user {
    background: #e8f0fb;
    border-radius: 10px;
    padding: 12px 14px;
    margin-bottom: 10px;
    font-size: 13px;
    color: #1a3c6e;
    font-weight: 600;
}
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────
# AUTH HELPERS
# ─────────────────────────────────────────

HARDCODED_API_KEY = "gsk_Yn3bFa7pUEdLkjMnoPqXrstUvWxYz"  # replace with your real key

def validate_credentials(user_id: str, password: str) -> bool:
    uid_ok  = bool(re.fullmatch(r"\d{5}", user_id))
    pwd_ok  = bool(re.fullmatch(r"VIG\d{4}", password))
    return uid_ok and pwd_ok

# ─────────────────────────────────────────
# SESSION INIT
# ─────────────────────────────────────────

if "logged_in" not in st.session_state:
    st.session_state["logged_in"] = False
if "faculty_id" not in st.session_state:
    st.session_state["faculty_id"] = ""
if "api_key" not in st.session_state:
    st.session_state["api_key"] = ""

# ─────────────────────────────────────────
# LOGIN PAGE
# ─────────────────────────────────────────

if not st.session_state["logged_in"]:

    _, col, _ = st.columns([1, 1.2, 1])

    with col:
        st.markdown("""
        <div class="login-card">
            <div class="login-logo">🎓</div>
            <div class="login-title">Vignan's Faculty Portal</div>
            <div class="login-sub">Question Paper Generator</div>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("<br>", unsafe_allow_html=True)

        with st.container():
            user_id  = st.text_input("Faculty User ID", placeholder="5-digit ID  e.g. 10234", max_chars=5)
            password = st.text_input("Password", type="password", placeholder="VIG followed by 4 digits  e.g. VIG1234")
            api_key_input = st.text_input("GROQ API Key", type="password", placeholder="gsk_...")

            st.markdown("<br>", unsafe_allow_html=True)
            login_btn = st.button("🔐  Login", use_container_width=True, type="primary")

            if login_btn:
                if not user_id or not password or not api_key_input:
                    st.error("All fields are required.")
                elif not validate_credentials(user_id, password):
                    if not re.fullmatch(r"\d{5}", user_id):
                        st.error("User ID must be exactly 5 digits.")
                    else:
                        st.error("Password must be in the format VIG followed by 4 digits (e.g. VIG1234).")
                else:
                    st.session_state["logged_in"]  = True
                    st.session_state["faculty_id"] = user_id
                    st.session_state["api_key"]    = api_key_input
                    st.rerun()

        st.markdown("""
        <div style='text-align:center;margin-top:18px;font-size:12px;color:#aaa;'>
            Vignan's University · AI-Powered Exam Tools
        </div>
        """, unsafe_allow_html=True)

    st.stop()

# ─────────────────────────────────────────
# ── AUTHENTICATED APP BELOW ──
# ─────────────────────────────────────────

api_key    = st.session_state["api_key"]
faculty_id = st.session_state["faculty_id"]

# ── TITLE BAR ──
st.markdown(f"""
<div class="title-bar">
    <div class="title-left">
        <div class="title-icon">🎓</div>
        <div class="title-text">
            <h1>Question Paper Generator</h1>
            <p>Vignan's University · AI-Powered · RAG · LangGraph</p>
        </div>
    </div>
    <div class="title-right">
        <span class="faculty-chip">Faculty ID: {faculty_id}</span>
    </div>
</div>
""", unsafe_allow_html=True)

# ── SIDEBAR ──
with st.sidebar:
    st.markdown(f'<div class="sidebar-user">👤 Faculty ID: {faculty_id}</div>', unsafe_allow_html=True)
    st.markdown("### 🗺️ Workflow")
    st.markdown("""
    1. **Exam Details** – type & difficulty  
    2. **Question Format** – sections & counts  
    3. **Knowledge Source** – topics or files  
    4. **Generate** – AI builds the paper  
    5. **Approve & Download** – export PDF
    """)
    st.markdown("---")
    if st.button("🚪 Logout", use_container_width=True):
        for k in ["logged_in", "faculty_id", "api_key", "paper", "exam_meta"]:
            st.session_state.pop(k, None)
        st.rerun()
    st.caption("Powered by Groq LLaMA · ChromaDB · LangGraph")

# ─────────────────────────────────────────
# EXAM META HELPERS
# ─────────────────────────────────────────

EXAM_TIME = {
    "T1":    "1 Hour",
    "T4":    "1 Hour 30 Minutes",
    "Final": "2 Hours 30 Minutes"
}
EXAM_MARKS = {
    "T1":    "20",
    "T4":    "20",
    "Final": "60"
}

def get_exam_pattern(exam_type: str) -> str:
    if "T1" in exam_type:
        return "T1"
    elif "T4" in exam_type:
        return "T4"
    return "Final"

# ─────────────────────────────────────────
# LLM
# ─────────────────────────────────────────

llm = ChatGroq(
    model_name="llama-3.1-8b-instant",
    temperature=0.3,
    api_key=api_key
)

# ─────────────────────────────────────────
# VECTOR DB
# ─────────────────────────────────────────

@st.cache_resource
def get_collection():
    client = chromadb.Client()
    return client.get_or_create_collection("syllabus_db")

collection     = get_collection()
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# ─────────────────────────────────────────
# FILE EXTRACTION
# ─────────────────────────────────────────

def extract_ppt_text(file):
    prs  = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_documents(files):
    docs = []
    for file in files:
        if file.name.endswith(".pdf"):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(file.read())
                tmp_path = tmp.name
            loader = PyPDFLoader(tmp_path)
            pages  = loader.load()
            docs.extend([p.page_content for p in pages])
            os.unlink(tmp_path)
        elif file.name.endswith(".pptx"):
            docs.append(extract_ppt_text(file))
    return docs

def process_documents(files):
    docs = load_documents(files)
    if not docs:
        return
    splitter   = RecursiveCharacterTextSplitter(chunk_size=400, chunk_overlap=50)
    chunks     = splitter.split_text("\n".join(docs))
    embeddings = embedding_model.encode(chunks)
    for i, chunk in enumerate(chunks):
        collection.add(
            documents  = [chunk],
            embeddings = [embeddings[i].tolist()],
            ids        = [f"doc_{i}_{abs(hash(chunk)) % 1000000}"]
        )

# ─────────────────────────────────────────
# LOGO DOWNLOAD
# ─────────────────────────────────────────

def get_vignan_logo() -> str | None:
    logo_url = "https://vignan.ac.in/newvignan/assets/images/logo.png"
    try:
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        urllib.request.urlretrieve(logo_url, tmp.name)
        return tmp.name
    except Exception:
        return None

# ─────────────────────────────────────────
# PDF GENERATOR
# ─────────────────────────────────────────

def generate_question_paper_pdf(text, exam_type, subject, branch, college):
    pattern   = get_exam_pattern(exam_type)
    time_str  = EXAM_TIME.get(pattern,  "2 Hours")
    marks_str = EXAM_MARKS.get(pattern, "20")

    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")

    doc = SimpleDocTemplate(
        tmp_file.name,
        pagesize=A4,
        rightMargin=50, leftMargin=50,
        topMargin=40,   bottomMargin=50
    )

    styles = getSampleStyleSheet()

    uni_style = ParagraphStyle(
        "uni", parent=styles["Heading1"],
        alignment=TA_CENTER, fontSize=13, spaceAfter=2,
        textColor=colors.HexColor("#1a1a2e"), fontName="Helvetica-Bold"
    )
    dept_style = ParagraphStyle(
        "dept", parent=styles["Normal"],
        alignment=TA_CENTER, fontSize=11, spaceAfter=2,
        textColor=colors.HexColor("#333")
    )
    section_style = ParagraphStyle(
        "section", parent=styles["Normal"],
        fontSize=12, spaceBefore=12, spaceAfter=5,
        textColor=colors.HexColor("#1a3c6e"), fontName="Helvetica-Bold",
        alignment=TA_CENTER
    )
    question_style = ParagraphStyle(
        "question", parent=styles["Normal"],
        fontSize=11, spaceAfter=5, leftIndent=8
    )
    subq_style = ParagraphStyle(
        "subq", parent=styles["Normal"],
        fontSize=11, spaceAfter=4, leftIndent=24
    )

    elements  = []
    logo_path = get_vignan_logo()

    # ── HEADER: Logo + University Name ──
    if logo_path:
        try:
            logo_img   = Image(logo_path, width=2.8*cm, height=2.8*cm)
            name_block = [
                Paragraph(college.upper(), uni_style),
                Paragraph("Department of ACSE", dept_style),
            ]
            header_tbl = Table([[logo_img, name_block]], colWidths=[3.2*cm, 13.0*cm])
            header_tbl.setStyle(TableStyle([
                ("VALIGN",      (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (0, 0),   0),
            ]))
            elements.append(header_tbl)
        except Exception:
            elements.append(Paragraph(college.upper(), uni_style))
            elements.append(Paragraph("Department of ACSE", dept_style))
    else:
        elements.append(Paragraph(college.upper(), uni_style))
        elements.append(Paragraph("Department of ACSE", dept_style))

    elements.append(Spacer(1, 6))
    elements.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor("#1a3c6e")))
    elements.append(Spacer(1, 4))

    # ── INFO TABLE ──
    info_data = [
        [
            Paragraph(f"<b>Subject:</b> {subject}", styles["Normal"]),
            Paragraph(f"<b>Exam:</b> {exam_type}", styles["Normal"]),
        ],
        [
            Paragraph(f"<b>Branch/Class:</b> {branch}", styles["Normal"]),
            Paragraph(
                f"<b>Time:</b> {time_str} &nbsp;&nbsp; <b>Max Marks:</b> {marks_str}",
                styles["Normal"]
            ),
        ],
    ]
    info_tbl = Table(info_data, colWidths=[8.5*cm, 8.5*cm])
    info_tbl.setStyle(TableStyle([
        ("FONTSIZE",      (0, 0), (-1, -1), 11),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING",    (0, 0), (-1, -1), 2),
    ]))
    elements.append(info_tbl)
    elements.append(Spacer(1, 4))
    elements.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#aaa")))
    elements.append(Spacer(1, 8))

    # ── QUESTION CONTENT ──
    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            elements.append(Spacer(1, 4))
            continue

        lower = stripped.lower()

        is_section = (
            lower.startswith("part a") or
            lower.startswith("part b") or
            lower.startswith("part c")
        )
        if is_section:
            elements.append(Spacer(1, 6))
            elements.append(Paragraph(stripped.upper(), section_style))
            elements.append(HRFlowable(
                width="55%", thickness=0.5,
                color=colors.HexColor("#1a3c6e"), hAlign="CENTER"
            ))
            elements.append(Spacer(1, 6))
            continue

        is_subq = (
            len(stripped) > 2 and
            stripped[0].isalpha() and
            stripped[1] in (")", ".") and
            stripped[0].islower()
        )
        if is_subq:
            elements.append(Paragraph(stripped, subq_style))
        else:
            elements.append(Paragraph(stripped, question_style))

    doc.build(elements)

    if logo_path:
        try:
            os.unlink(logo_path)
        except Exception:
            pass

    return tmp_file.name

# ─────────────────────────────────────────
# LANGGRAPH STATE
# ─────────────────────────────────────────

class PaperState(TypedDict):
    query:          str
    context:        str
    topics:         str
    section_config: dict
    exam_type:      str
    difficulty:     str
    questions:      str
    response:       str
    validated:      bool   # True = content is academic, False = rejected

# ─────────────────────────────────────────
# GRAPH NODES
# ─────────────────────────────────────────

def retrieval_agent(state: PaperState) -> PaperState:
    query = state["topics"].strip()

    # ── No input at all ──
    if not query:
        return {**state, "context": "", "validated": False, "topics": "INVALID_DATA"}

    # ── No documents uploaded → knowledge base is empty ──
    if collection.count() == 0:
        return {**state, "context": "", "validated": False, "topics": "NO_KB"}

    # ── Query the knowledge base and check similarity scores ──
    results = collection.query(
        query_texts=[query],
        n_results=5,
        include=["documents", "distances"]
    )

    docs      = results["documents"][0]      # list of retrieved chunks
    distances = results["distances"][0]      # L2 distances (lower = more similar)

    # ChromaDB returns L2 distances; convert to a 0-1 relevance score
    # We reject if the BEST (lowest) distance is too high → content not in KB
    DISTANCE_THRESHOLD = 1.2   # tuned: < 1.2 means meaningful match, > 1.2 = unrelated

    best_distance = min(distances) if distances else 999
    if best_distance > DISTANCE_THRESHOLD:
        return {**state, "context": "", "validated": False, "topics": "NOT_IN_KB"}

    # ── Only use chunks that are actually relevant ──
    relevant_chunks = [
        doc for doc, dist in zip(docs, distances)
        if dist <= DISTANCE_THRESHOLD
    ]
    context = "\n".join(relevant_chunks)

    return {**state, "context": context, "validated": True}


def validator_agent(state: PaperState) -> PaperState:
    """Pass-through — all validation is already done in retrieval_agent via KB similarity."""
    return state



def route_after_validation(state: PaperState) -> str:
    """Conditional edge: go to generator only if content passed validation."""
    if state["validated"]:
        return "generator"
    return "formatter"   # skip generator, go straight to formatter with error message


def generator_agent(state: PaperState) -> PaperState:
    # If validation failed, set an error message and pass through
    if not state.get("validated", True):
        return {**state, "questions": "INVALID_CONTENT"}

    topics     = state["topics"]
    context    = state["context"]
    config     = state["section_config"]
    exam_type  = state["exam_type"]
    difficulty = state["difficulty"]

    pattern = config.get("pattern", "T1")

    def sub_labels(n):
        return list("abcde")[:n]

    # ── T1 ──
    if pattern == "T1":
        pA_q, pA_sub = config.get("partA_q", 2), config.get("partA_sub", 2)
        pB_q, pB_sub = config.get("partB_q", 2), config.get("partB_sub", 2)

        partA_scaffold = []
        for q in range(1, pA_q + 1):
            partA_scaffold.append(f"Q{q}.")
            for s in sub_labels(pA_sub):
                partA_scaffold.append(f"  {s}) <fill question>")

        partB_scaffold = []
        for q in range(1, pB_q + 1):
            partB_scaffold.append(f"Q{q}.")
            for s in sub_labels(pB_sub):
                partB_scaffold.append(f"  {s}) <fill question>")

        format_block = (
            "PART A\n" + "\n".join(partA_scaffold) +
            "\n\nPART B\n" + "\n".join(partB_scaffold)
        )
        rules = (
            f"- PART A: exactly {pA_q} questions, each with exactly {pA_sub} sub-parts.\n"
            f"- PART B: exactly {pB_q} questions, each with exactly {pB_sub} sub-parts.\n"
            "- PART A questions should be short/conceptual.\n"
            "- PART B questions should be analytical/descriptive.\n"
        )

    # ── T4 ──
    elif pattern == "T4":
        mcq_count        = config.get("mcq", 10)
        pB_q, pB_sub     = config.get("partB_q", 3), config.get("partB_sub", 2)

        mcq_scaffold = []
        for q in range(1, mcq_count + 1):
            mcq_scaffold.append(
                f"Q{q}. <fill MCQ question>\n"
                f"  A) <option>  B) <option>  C) <option>  D) <option>\n"
                f"  Ans: <correct letter>"
            )

        partB_scaffold = []
        for q in range(1, pB_q + 1):
            partB_scaffold.append(f"Q{q}.")
            for s in sub_labels(pB_sub):
                partB_scaffold.append(f"  {s}) <fill question>")

        format_block = (
            "PART A\n" + "\n".join(mcq_scaffold) +
            "\n\nPART B\n" + "\n".join(partB_scaffold)
        )
        rules = (
            f"- PART A: exactly {mcq_count} MCQs, each with 4 options (A B C D) and Ans line.\n"
            f"- PART B: exactly {pB_q} questions, each with exactly {pB_sub} sub-parts.\n"
        )

    # ── FINAL SEMESTER ── (Part A and Part B only)
    else:
        pA_q, pA_sub = config.get("partA_q", 5), config.get("partA_sub", 2)
        pB_q, pB_sub = config.get("partB_q", 4), config.get("partB_sub", 2)

        partA_scaffold = []
        for q in range(1, pA_q + 1):
            partA_scaffold.append(f"Q{q}.")
            for s in sub_labels(pA_sub):
                partA_scaffold.append(f"  {s}) <fill question>")

        partB_scaffold = []
        for q in range(1, pB_q + 1):
            partB_scaffold.append(f"Q{q}.")
            for s in sub_labels(pB_sub):
                partB_scaffold.append(f"  {s}) <fill question>")

        format_block = (
            "PART A\n" + "\n".join(partA_scaffold) +
            "\n\nPART B\n" + "\n".join(partB_scaffold)
        )
        rules = (
            f"- PART A: exactly {pA_q} questions, each with exactly {pA_sub} sub-parts (short/conceptual level).\n"
            f"- PART B: exactly {pB_q} questions, each with exactly {pB_sub} sub-parts (detailed/analytical level).\n"
        )

    prompt = f"""You are a university professor generating a formal exam question paper.

Syllabus Context:
{context}

Topics:
{topics}

Exam Type: {exam_type}
Difficulty: {difficulty}

TASK: Fill in every <fill question>, <fill MCQ question>, <option>, and <correct letter> placeholder in the skeleton below with real, well-formed questions based on the topics.

{format_block}

STRICT RULES (follow exactly):
{rules}- Output ONLY the filled skeleton — nothing before PART A, nothing after the last question.
- Do NOT change the structure, numbering, or labels.
- Do NOT add any section title other than PART A and PART B.
- Do NOT add marks, time, instructions, introductions, or closing remarks.
"""

    result = llm.invoke(prompt)
    return {**state, "questions": result.content}


def formatting_agent(state: PaperState) -> PaperState:
    q      = state.get("questions", "")
    topics = state.get("topics", "")

    if q == "INVALID_CONTENT" or not state.get("validated", True):
        if topics == "NO_KB":
            msg = "NO_KB"
        elif topics == "NOT_IN_KB":
            msg = "NOT_IN_KB"
        else:
            msg = "INVALID_CONTENT"
        return {**state, "response": msg}

    return {**state, "response": q}

# ─────────────────────────────────────────
# BUILD GRAPH
# ─────────────────────────────────────────

graph = StateGraph(PaperState)
graph.add_node("retrieval",  retrieval_agent)
graph.add_node("validator",  validator_agent)
graph.add_node("generator",  generator_agent)
graph.add_node("formatter",  formatting_agent)
graph.add_edge(START, "retrieval")
graph.add_edge("retrieval", "validator")
graph.add_conditional_edges("validator", route_after_validation, {
    "generator": "generator",
    "formatter": "formatter",
})
graph.add_edge("generator", "formatter")
graph.add_edge("formatter", END)
workflow = graph.compile()

# ─────────────────────────────────────────
# STEP 1 — EXAM DETAILS
# ─────────────────────────────────────────

st.markdown('<div class="card">', unsafe_allow_html=True)
st.markdown('<p class="step-header">📋 Step 1 · Exam Details</p>', unsafe_allow_html=True)

col1, col2, col3, col4 = st.columns(4)
with col1:
    exam_type = st.selectbox(
        "Exam Type",
        ["Mid1-T1", "Mid1-T4", "Mid2-T1", "Mid2-T4", "Final Semester"],
        index=None, placeholder="Select..."
    )
with col2:
    difficulty = st.selectbox(
        "Difficulty",
        ["Easy", "Moderate", "Difficult"],
        index=None, placeholder="Select..."
    )
with col3:
    subject = st.text_input("Subject Name", placeholder="e.g. Computer Networks")
with col4:
    branch  = st.text_input("Branch / Class", placeholder="e.g. III AIML B")

college = st.text_input(
    "University / College Name",
    value="Vignan's Foundation for Science, Technology and Research"
)

if exam_type:
    pat = get_exam_pattern(exam_type)
    st.info(f"⏱️  Duration: **{EXAM_TIME[pat]}**   |   📝  Max Marks: **{EXAM_MARKS[pat]}**")

st.markdown('</div>', unsafe_allow_html=True)

# ─────────────────────────────────────────
# STEP 2 — QUESTION FORMAT
# ─────────────────────────────────────────

section_config = {}

if exam_type and difficulty:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.markdown('<p class="step-header">🗂️ Step 2 · Configure Question Format</p>', unsafe_allow_html=True)

    if "T1" in exam_type:
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("**Part A**")
            partA_q   = st.number_input("Number of Questions (Part A)",        1, 10, value=2)
            partA_sub = st.number_input("Sub-questions per Question (Part A)",  1, 5,  value=2)
        with col2:
            st.markdown("**Part B**")
            partB_q   = st.number_input("Number of Questions (Part B)",        1, 10, value=2)
            partB_sub = st.number_input("Sub-questions per Question (Part B)", 1, 5,  value=2)
        section_config = {
            "pattern": "T1",
            "partA_q": partA_q, "partA_sub": partA_sub,
            "partB_q": partB_q, "partB_sub": partB_sub
        }

    elif "T4" in exam_type:
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("**Part A – MCQs**")
            mcq = st.number_input("Number of MCQs", 1, 50, value=10)
        with col2:
            st.markdown("**Part B**")
            partB_q   = st.number_input("Number of Questions (Part B)",        1, 10, value=3)
            partB_sub = st.number_input("Sub-questions per Question (Part B)", 1, 5,  value=2)
        section_config = {
            "pattern": "T4",
            "mcq": mcq,
            "partB_q": partB_q, "partB_sub": partB_sub
        }

    elif "Final" in exam_type:
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("**Part A**")
            partA_q   = st.number_input("Number of Questions (Part A)",        1, 20, value=5)
            partA_sub = st.number_input("Sub-questions per Question (Part A)", 1, 5,  value=2)
        with col2:
            st.markdown("**Part B**")
            partB_q   = st.number_input("Number of Questions (Part B)",        1, 10, value=4)
            partB_sub = st.number_input("Sub-questions per Question (Part B)", 1, 5,  value=2)
        section_config = {
            "pattern":  "Final",
            "partA_q":  partA_q,  "partA_sub": partA_sub,
            "partB_q":  partB_q,  "partB_sub": partB_sub,
        }

    st.markdown('</div>', unsafe_allow_html=True)

# ─────────────────────────────────────────
# STEP 3 — KNOWLEDGE SOURCE
# ─────────────────────────────────────────

st.markdown('<div class="card">', unsafe_allow_html=True)
st.markdown('<p class="step-header">📂 Step 3 · Provide Knowledge Source</p>', unsafe_allow_html=True)

tab1, tab2 = st.tabs(["✏️ Manual Topics", "📁 Upload Files"])

with tab1:
    manual_topics = st.text_area(
        "Enter syllabus topics / key concepts",
        placeholder="e.g. OSI Model, TCP/IP, Routing Protocols, DNS, HTTP, Subnetting...",
        height=120
    )

with tab2:
    files = st.file_uploader(
        "Upload PDF or PPTX files",
        type=["pdf", "pptx"],
        accept_multiple_files=True
    )
    if files:
        with st.spinner("Indexing documents..."):
            process_documents(files)
        st.success(f"✅ {len(files)} file(s) indexed successfully.")

st.markdown('</div>', unsafe_allow_html=True)

# ─────────────────────────────────────────
# STEP 4 — GENERATE
# ─────────────────────────────────────────

col_l, col_c, col_r = st.columns([1, 2, 1])
with col_c:
    generate = st.button("🚀 Generate Question Paper", use_container_width=True, type="primary")

if generate:
    if not exam_type or not difficulty:
        st.error("Please complete Step 1 — select exam type and difficulty.")
        st.stop()
    if not section_config:
        st.error("Please complete Step 2 — configure question format.")
        st.stop()
    if not manual_topics.strip():
        st.error("Please enter topics in Step 3 to search your knowledge base.")
        st.stop()
    if collection.count() == 0:
        st.error("Please upload a syllabus or document in Step 3 before generating.")
        st.stop()

    with st.spinner("🔍 Validating content and generating questions..."):
        initial_state: PaperState = {
            "query":          manual_topics,
            "context":        "",
            "topics":         manual_topics,
            "section_config": section_config,
            "exam_type":      exam_type,
            "difficulty":     difficulty,
            "questions":      "",
            "response":       "",
            "validated":      True
        }
        result   = workflow.invoke(initial_state)
        response = result["response"]

        ERROR_MESSAGES = {
            "NO_KB": (
                "❌ **No knowledge base found.**\n\n"
                "You have not uploaded any syllabus, notes, or document yet. "
                "Please go to **Step 3 → Upload Files** and upload a relevant PDF or PPTX file first."
            ),
            "NOT_IN_KB": (
                "❌ **Content not found in your knowledge base.**\n\n"
                "The topics you entered do not match anything in your uploaded documents. "
                "Please enter topics that are actually covered in your uploaded syllabus/notes, "
                "or upload the correct document."
            ),
            "INVALID_CONTENT": (
                "❌ **Invalid or unrelated content detected.**\n\n"
                "The input does not appear to be academic/educational material. "
                "Please enter valid syllabus topics or upload a relevant subject document."
            ),
        }

        if response in ERROR_MESSAGES:
            st.session_state.pop("paper", None)
            st.error(ERROR_MESSAGES[response])
        else:
            st.session_state["paper"] = response
            st.session_state["exam_meta"] = {
                "exam_type": exam_type,
                "subject":   subject or "Subject",
                "branch":    branch  or "Branch",
                "college":   college or "University"
            }

# ─────────────────────────────────────────
# STEP 5 — REVIEW & DOWNLOAD
# ─────────────────────────────────────────

if "paper" in st.session_state:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.markdown('<p class="step-header">📄 Step 4 · Review & Edit Question Paper</p>', unsafe_allow_html=True)

    st.session_state["paper"] = st.text_area(
        "You can edit the question paper below before approving:",
        value=st.session_state["paper"],
        height=450
    )
    st.markdown('</div>', unsafe_allow_html=True)

    col_l, col_c, col_r = st.columns([1, 2, 1])
    with col_c:
        approve = st.button("✅ Approve & Download PDF", use_container_width=True)

    if approve:
        meta = st.session_state.get("exam_meta", {})
        with st.spinner("📄 Building PDF..."):
            pdf_path = generate_question_paper_pdf(
                text      = st.session_state["paper"],
                exam_type = meta.get("exam_type", "Exam"),
                subject   = meta.get("subject",   "Subject"),
                branch    = meta.get("branch",     "Branch"),
                college   = meta.get("college",    "University")
            )

        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()

        fname = "{}_{}_{}.pdf".format(
            meta.get("subject",   "QuestionPaper").replace(" ", "_"),
            meta.get("branch",    "").replace(" ", "_"),
            meta.get("exam_type", "Exam").replace(" ", "_")
        )

        st.download_button(
            label               = "⬇️ Download Question Paper PDF",
            data                = pdf_bytes,
            file_name           = fname,
            mime                = "application/pdf",
            use_container_width = True
        )
        st.success("✅ PDF ready! Click the button above to download.")
        os.unlink(pdf_path)